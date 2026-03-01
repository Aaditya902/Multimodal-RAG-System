from typing import List, Dict, Any
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
from openpyxl.utils import get_column_letter
import csv
import streamlit as st
import os

from .base import BaseProcessor
from models.chunk import Chunk

class DocumentProcessor(BaseProcessor):

    def __init__(self):
        self._extensions = ['.docx', '.doc', '.xlsx', '.xls', '.csv', '.pptx', '.ppt']
    
    def process(self, file_path: str, source_name: str) -> List[Chunk]:
        
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext in ['.docx', '.doc']:
            return self._process_word(file_path, source_name)
        elif ext in ['.xlsx', '.xls', '.csv']:
            return self._process_excel_enhanced(file_path, source_name)
        elif ext in ['.pptx', '.ppt']:
            return self._process_powerpoint(file_path, source_name)
        else:
            return []
    
    def _process_excel_enhanced(self, file_path: str, source_name: str) -> List[Chunk]:
        
        chunks = []
        
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
                sheet_name = "CSV Data"
            else:
                xl = pd.ExcelFile(file_path)
                sheet_names = xl.sheet_names
                
                for sheet_idx, sheet_name in enumerate(sheet_names):
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    sheet_chunks = self._process_dataframe(df, source_name, sheet_name, sheet_idx)
                    chunks.extend(sheet_chunks)
                
                return chunks
            
            chunks = self._process_dataframe(df, source_name, "CSV Data", 0)
            
        except Exception as e:
            st.warning(f"Excel processing failed: {str(e)}")
            chunks = self._process_excel_basic(file_path, source_name)
        
        return chunks
    
    def _process_dataframe(self, df: pd.DataFrame, source_name: str, 
                          sheet_name: str, sheet_idx: int) -> List[Chunk]:
        """Process a pandas DataFrame into multiple chunks"""
        
        chunks = []
        
        summary = self._create_excel_summary(df, sheet_name)
        chunks.append(Chunk(
            content=summary,
            source_file=source_name,
            file_type='excel_summary',
            chunk_index=len(chunks),
            metadata={
                'sheet': sheet_name,
                'type': 'summary',
                'rows': len(df),
                'columns': len(df.columns)
            }
        ))
        
        col_info = self._create_column_info(df)
        chunks.append(Chunk(
            content=col_info,
            source_file=source_name,
            file_type='excel_columns',
            chunk_index=len(chunks),
            metadata={
                'sheet': sheet_name,
                'type': 'columns'
            }
        ))
        
        data_chunks = self._create_data_chunks(df, source_name, sheet_name)
        chunks.extend(data_chunks)
        
        stats = self._create_statistical_summary(df)
        if stats:
            chunks.append(Chunk(
                content=stats,
                source_file=source_name,
                file_type='excel_statistics',
                chunk_index=len(chunks),
                metadata={
                    'sheet': sheet_name,
                    'type': 'statistics'
                }
            ))
        
        return chunks
    
    def _create_excel_summary(self, df: pd.DataFrame, sheet_name: str) -> str:
        """Create a summary of the Excel data"""
        
        summary = []
        summary.append(f"=== Excel Sheet: {sheet_name} ===")
        summary.append(f"Total Rows: {len(df)}")
        summary.append(f"Total Columns: {len(df.columns)}")
        summary.append(f"Column Names: {', '.join(df.columns)}")
        
        dtypes = df.dtypes.value_counts()
        summary.append("\nData Types:")
        for dtype, count in dtypes.items():
            summary.append(f"  - {dtype}: {count} columns")
        
        missing = df.isnull().sum().sum()
        if missing > 0:
            summary.append(f"\nMissing Values: {missing}")
        
        summary.append("\nSample Data (First 3 rows):")
        sample = df.head(3).to_string()
        summary.append(sample)
        
        return '\n'.join(summary)
    
    def _create_column_info(self, df: pd.DataFrame) -> str:
        """Create detailed column information"""
        
        col_info = ["=== Column Details ==="]
        
        for col in df.columns:
            info = []
            info.append(f"\nColumn: {col}")
            info.append(f"  Type: {df[col].dtype}")
            info.append(f"  Unique Values: {df[col].nunique()}")
            
            if df[col].dtype in ['int64', 'float64']:
                info.append(f"  Range: {df[col].min()} to {df[col].max()}")
                info.append(f"  Mean: {df[col].mean():.2f}")
            
            sample_vals = df[col].dropna().head(3).tolist()
            if sample_vals:
                info.append(f"  Sample: {sample_vals}")
            
            col_info.extend(info)
        
        return '\n'.join(col_info)
    
    def _create_data_chunks(self, df: pd.DataFrame, source_name: str, 
                           sheet_name: str, chunk_size: int = 50) -> List[Chunk]:
        """Split dataframe into manageable chunks"""
        
        chunks = []
        total_rows = len(df)
        
        for start_row in range(0, total_rows, chunk_size):
            end_row = min(start_row + chunk_size, total_rows)
            chunk_df = df.iloc[start_row:end_row]
            
            if len(chunk_df) <= 20:
                data_str = chunk_df.to_string()
            else:
                data_str = f"Rows {start_row+1}-{end_row} of {total_rows}:\n"
                data_str += chunk_df.head(10).to_string()
                data_str += f"\n... and {len(chunk_df)-10} more rows"
            
            chunk = Chunk(
                content=f"[Excel Data - {sheet_name} - Rows {start_row+1}-{end_row}]\n{data_str}",
                source_file=source_name,
                file_type='excel_data',
                chunk_index=len(chunks),
                metadata={
                    'sheet': sheet_name,
                    'start_row': start_row,
                    'end_row': end_row,
                    'rows': len(chunk_df)
                }
            )
            chunks.append(chunk)
        
        return chunks
    
    def _create_statistical_summary(self, df: pd.DataFrame) -> str:
        """Create statistical summary for numeric columns"""
        
        numeric_cols = df.select_dtypes(include=['int64', 'float64']).columns
        
        if len(numeric_cols) == 0:
            return ""
        
        stats = ["=== Statistical Summary (Numeric Columns) ==="]
        
        for col in numeric_cols:
            stats.append(f"\n{col}:")
            stats.append(f"  Count: {df[col].count()}")
            stats.append(f"  Mean: {df[col].mean():.2f}")
            stats.append(f"  Std: {df[col].std():.2f}")
            stats.append(f"  Min: {df[col].min()}")
            stats.append(f"  25%: {df[col].quantile(0.25)}")
            stats.append(f"  50%: {df[col].quantile(0.50)}")
            stats.append(f"  75%: {df[col].quantile(0.75)}")
            stats.append(f"  Max: {df[col].max()}")
        
        return '\n'.join(stats)
    
    def _process_excel_basic(self, file_path: str, source_name: str) -> List[Chunk]:
        """Fallback basic Excel processing"""
        chunks = []
        
        try:
            if file_path.endswith('.csv'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    rows = list(reader)
                    text = "CSV Data:\n"
                    for row in rows[:100]:  
                        text += ' | '.join(row) + '\n'
            else:
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text = "Excel Data:\n"
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    text += f"\n--- Sheet: {sheet} ---\n"
                    for row in ws.iter_rows(max_row=50, values_only=True):
                        if any(cell for cell in row):
                            text += ' | '.join(str(cell) for cell in row if cell) + '\n'
            
            chunk = Chunk(
                content=text,
                source_file=source_name,
                file_type='excel_basic',
                chunk_index=0,
                metadata={'type': 'basic_fallback'}
            )
            chunks.append(chunk)
            
        except Exception as e:
            st.warning(f"Basic Excel processing failed: {str(e)}")
        
        return chunks
    
    def _process_word(self, file_path: str, source_name: str) -> List[Chunk]:
        """Process Word document (existing code)"""
        chunks = []
        
        try:
            doc = DocxDocument(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            for table in doc.tables:
                table_text = ["Table:"]
                for row in table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells if cell.text)
                    if row_text:
                        table_text.append(row_text)
                if len(table_text) > 1:
                    text_parts.append('\n'.join(table_text))
            
            if text_parts:
                full_text = '\n\n'.join(text_parts)
                chunk_texts = self._chunk_text(full_text)
                
                for i, chunk_text in enumerate(chunk_texts):
                    chunk = Chunk(
                        content=chunk_text,
                        source_file=source_name,
                        file_type='word',
                        chunk_index=i,
                        metadata={'type': 'word_document'}
                    )
                    chunks.append(chunk)
        
        except Exception as e:
            st.warning(f"Word processing failed: {str(e)}")
        
        return chunks
    
    def _process_powerpoint(self, file_path: str, source_name: str) -> List[Chunk]:
        """Process PowerPoint presentation"""
        chunks = []
        
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"\n=== Slide {slide_num + 1} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    if hasattr(shape, "table"):
                        table_text = ["Table:"]
                        for row in shape.table.rows:
                            row_text = ' | '.join(cell.text for cell in row.cells if cell.text)
                            if row_text:
                                table_text.append(row_text)
                        if len(table_text) > 1:
                            slide_text.append('\n'.join(table_text))
                
                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))
            
            if text_parts:
                full_text = '\n'.join(text_parts)
                chunk_texts = self._chunk_text(full_text)
                
                for i, chunk_text in enumerate(chunk_texts):
                    chunk = Chunk(
                        content=chunk_text,
                        source_file=source_name,
                        file_type='powerpoint',
                        chunk_index=i,
                        metadata={'type': 'presentation'}
                    )
                    chunks.append(chunk)
        
        except Exception as e:
            st.warning(f"PowerPoint processing failed: {str(e)}")
        
        return chunks
    
    def _chunk_text(self, text: str, max_size: int = 1000) -> List[str]:
        """Split text into chunks"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
        
        for word in words:
            word_length = len(word) + 1
            if current_length + word_length > max_size and current_chunk:
                chunks.append(' '.join(current_chunk))
                current_chunk = []
                current_length = 0
            current_chunk.append(word)
            current_length += word_length
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    def supports(self, file_type: str) -> bool:
        return file_type.lower() in ['word', 'excel', 'powerpoint']
    
    @property
    def supported_extensions(self) -> List[str]:
        return self._extensions